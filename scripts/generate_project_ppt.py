from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
OUTPUT = DOCS_DIR / "10-exam-integrity-system-presentation.pptx"


BG = RGBColor(245, 248, 252)
NAVY = RGBColor(16, 36, 62)
BLUE = RGBColor(26, 92, 171)
LIGHT = RGBColor(225, 236, 248)
GREEN = RGBColor(223, 242, 230)
TEXT = RGBColor(35, 43, 53)
MUTED = RGBColor(96, 110, 128)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(236, 106, 84)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.3), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.95), Inches(8.2), Inches(0.45))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(2.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_bullets(slide, items, left=0.8, top=1.6, width=5.0, height=4.8, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.bullet = True
        p.space_after = Pt(8)
    return box


def add_box(slide, left, top, width, height, text, fill_rgb=LIGHT, font_size=16, bold=True):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = NAVY
    return shape


def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = BLUE
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def add_two_col_table_like(slide, title_left, title_right, left_items, right_items):
    add_box(slide, 0.75, 1.55, 3.95, 0.45, title_left, NAVY, 16)
    add_box(slide, 5.0, 1.55, 3.95, 0.45, title_right, NAVY, 16)
    slide.shapes[-2].text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    slide.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    add_bullets(slide, left_items, left=0.9, top=2.1, width=3.6, height=4.7, font_size=17)
    add_bullets(slide, right_items, left=5.15, top=2.1, width=3.55, height=4.7, font_size=17)


def add_footer(slide, text="Exam Integrity System"):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(2.8), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)


def new_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title, subtitle)
    add_footer(slide)
    return slide


# Slide 1
slide = new_slide("Exam Integrity System", "Role-based online examination platform with integrity monitoring and audit-ready reporting")
hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.7), Inches(8.5), Inches(4.6))
hero.fill.solid()
hero.fill.fore_color.rgb = WHITE
hero.line.color.rgb = LIGHT
tf = hero.text_frame
tf.clear()
for idx, line in enumerate([
    "Project Focus",
    "Secure exam delivery, suspicious activity tracking, manual proctor review, controlled result publication, and cloud-backed audit reports."
]):
    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = line
    r.font.name = "Aptos"
    r.font.bold = idx == 0
    r.font.size = Pt(24 if idx == 0 else 20)
    r.font.color.rgb = NAVY if idx == 0 else TEXT
    p.space_after = Pt(20 if idx == 0 else 0)
for text, x in [("React + Vite", 1.2), ("Node + Express", 3.4), ("Neon Postgres", 5.75), ("EC2 + S3", 7.7)]:
    add_box(slide, x, 5.2, 1.15, 0.55, text, GREEN, 11)


# Slide 2
slide = new_slide("Problem Statement", "Why this project was needed")
add_bullets(slide, [
    "Basic online exam portals mostly stop at question delivery and answer submission.",
    "They often lack integrity monitoring, investigation workflow, and post-exam traceability.",
    "Institutions need role separation between admin, evaluator, proctor, and auditor.",
    "Generated reports and evidence should be stored securely and remain accessible for review."
], left=0.85, top=1.7, width=5.2, height=4.8, font_size=20)
add_box(slide, 6.45, 1.95, 2.0, 0.8, "Need", ACCENT, 22)
add_box(slide, 6.2, 3.0, 2.5, 2.2, "A complete exam lifecycle system\n\nnot just a quiz portal", LIGHT, 20)


# Slide 3
slide = new_slide("Solution Overview", "What the system delivers")
add_bullets(slide, [
    "Role-based platform for Admin, Student, Proctor, Evaluator, and Auditor.",
    "Student registration with email verification and scheduled exam access.",
    "Suspicious activity logging plus manual penalty and case review.",
    "Evaluator marking and admin-controlled result publication.",
    "S3-backed result reports and integrity evidence for audit review."
], left=0.8, top=1.75, width=5.2, height=4.9, font_size=19)
for y, txt, fill in [(1.9, "Exam Setup", LIGHT), (2.9, "Exam Attempt", GREEN), (3.9, "Integrity Review", LIGHT), (4.9, "Publish + Audit", GREEN)]:
    add_box(slide, 6.4, y, 2.2, 0.65, txt, fill, 18)
for a, b in [(2.55, 3.0), (3.55, 4.0), (4.55, 5.0)]:
    add_connector(slide, 7.5, a, 7.5, b)


# Slide 4
slide = new_slide("Role-Based Architecture", "Who does what in the system")
roles = [
    ("Admin", 0.8, 2.0, LIGHT),
    ("Student", 2.8, 2.0, GREEN),
    ("Proctor", 4.8, 2.0, LIGHT),
    ("Evaluator", 6.8, 2.0, GREEN),
    ("Auditor", 3.8, 4.4, LIGHT),
]
positions = {}
for name, x, y, fill in roles:
    positions[name] = (x + 0.7, y + 0.35)
    add_box(slide, x, y, 1.5, 0.7, name, fill, 18)
for start, end in [("Admin", "Student"), ("Student", "Proctor"), ("Proctor", "Evaluator"), ("Admin", "Auditor"), ("Evaluator", "Auditor")]:
    x1, y1 = positions[start]
    x2, y2 = positions[end]
    add_connector(slide, x1, y1, x2, y2)
add_bullets(slide, [
    "Admin creates exams and publishes final results.",
    "Student attempts the exam after verification.",
    "Proctor reviews suspicious events and case decisions.",
    "Evaluator awards marks.",
    "Auditor verifies logs, reports, and hash status."
], left=0.95, top=5.5, width=8.0, height=1.2, font_size=16)


# Slide 5
slide = new_slide("Technology Stack", "Main technologies used across the project")
add_two_col_table_like(
    slide,
    "Application Layers",
    "Cloud and Services",
    [
        "Frontend: React, Vite, JavaScript, Custom CSS",
        "Backend: Node.js, Express.js, REST APIs",
        "Database: PostgreSQL on Neon",
        "Database extensions: pgcrypto, citext",
    ],
    [
        "Hosting: AWS EC2",
        "Storage: AWS S3",
        "Reverse proxy: Nginx",
        "Process manager: PM2",
        "Email: SMTP",
        "Version control: GitHub",
    ],
)


# Slide 6
slide = new_slide("Dependencies Installed", "Key packages used in frontend and backend")
add_two_col_table_like(
    slide,
    "Backend Packages",
    "Frontend Packages",
    [
        "express, pg, nodemailer",
        "multer, cors, dotenv",
        "helmet, morgan, nodemon",
        "@aws-sdk/client-s3",
        "@aws-sdk/s3-request-presigner",
    ],
    [
        "react, react-dom",
        "vite, @vitejs/plugin-react",
        "eslint and react hooks plugin",
        "@eslint/js, globals",
        "@types/react, @types/react-dom",
    ],
)


# Slide 7
slide = new_slide("System Flow Diagram", "High-level project workflow from setup to audit")
steps = [
    ("Admin Setup", 0.7, 1.9, LIGHT),
    ("Student Login\nand Verification", 2.7, 1.9, GREEN),
    ("Exam Attempt\nwith Autosave", 4.9, 1.9, LIGHT),
    ("Integrity Event\nLogging", 7.1, 1.9, GREEN),
    ("Evaluation", 1.7, 4.4, GREEN),
    ("Proctor Review", 4.0, 4.4, LIGHT),
    ("Publish Results", 6.2, 4.4, GREEN),
]
anchors = []
for text, x, y, fill in steps:
    add_box(slide, x, y, 1.85, 0.85, text, fill, 16)
    anchors.append((x + 0.92, y + 0.42))
for start, end in [(0,1),(1,2),(2,3),(3,5),(2,4),(4,5),(5,6)]:
    x1, y1 = anchors[start]
    x2, y2 = anchors[end]
    add_connector(slide, x1, y1, x2, y2)


# Slide 8
slide = new_slide("Integrity and Result Flow", "How monitoring and publishing are controlled")
add_box(slide, 0.8, 2.0, 1.8, 0.75, "Suspicious Event", LIGHT, 16)
add_box(slide, 3.0, 2.0, 1.8, 0.75, "Penalty Assigned", GREEN, 16)
add_box(slide, 5.2, 2.0, 1.8, 0.75, "Case Opened", LIGHT, 16)
add_box(slide, 7.4, 2.0, 1.8, 0.75, "Decision Saved", GREEN, 16)
for x1, x2 in [(2.6,3.0),(4.8,5.2),(7.0,7.4)]:
    add_connector(slide, x1, 2.37, x2, 2.37)
add_box(slide, 1.6, 4.4, 2.3, 0.85, "Evaluator marks\nsubmission", LIGHT, 18)
add_box(slide, 4.2, 4.4, 2.3, 0.85, "Admin publishes\nresults", GREEN, 18)
add_box(slide, 6.8, 4.4, 2.0, 0.85, "S3 reports +\nresult emails", LIGHT, 18)
add_connector(slide, 3.9, 4.82, 4.2, 4.82)
add_connector(slide, 6.5, 4.82, 6.8, 4.82)
note = slide.shapes.add_textbox(Inches(0.95), Inches(5.8), Inches(8.1), Inches(0.7))
p = note.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Publishing is allowed only after all assigned students submit, all submissions are evaluated, and all opened integrity cases have decisions."
r.font.name = "Aptos"
r.font.size = Pt(16)
r.font.color.rgb = TEXT


# Slide 9
slide = new_slide("Deployment Architecture", "How frontend, backend, database, mail, and storage connect")
add_box(slide, 0.7, 2.2, 1.5, 0.8, "Browser", GREEN, 18)
add_box(slide, 2.6, 2.2, 1.8, 0.8, "Nginx on EC2", LIGHT, 18)
add_box(slide, 4.9, 2.2, 1.8, 0.8, "Express Backend", GREEN, 18)
add_box(slide, 7.2, 1.5, 1.6, 0.8, "Neon DB", LIGHT, 18)
add_box(slide, 7.2, 2.8, 1.6, 0.8, "SMTP Mail", LIGHT, 18)
add_box(slide, 7.2, 4.1, 1.6, 0.8, "AWS S3", GREEN, 18)
for pair in [((2.2,2.6),(2.6,2.6)), ((4.4,2.6),(4.9,2.6)), ((6.7,2.55),(7.2,1.9)), ((6.7,2.75),(7.2,3.2)), ((6.7,2.95),(7.2,4.5))]:
    add_connector(slide, pair[0][0], pair[0][1], pair[1][0], pair[1][1])
add_bullets(slide, [
    "Frontend build is served through Nginx.",
    "Backend runs on EC2 under PM2.",
    "Neon stores relational application data.",
    "SMTP handles verification and result emails.",
    "S3 stores result reports and integrity evidence.",
], left=0.9, top=5.7, width=5.8, height=1.0, font_size=15)


# Slide 10
slide = new_slide("Installation and Setup Steps", "From local machine to working application")
add_bullets(slide, [
    "Clone the repository and configure `.env` values.",
    "Install backend dependencies with `npm install` inside `backend`.",
    "Install frontend dependencies with `npm install` inside `frontend`.",
    "Create the Neon database and apply SQL schema scripts.",
    "Start backend on port 4000 and frontend on Vite dev server.",
    "Configure SMTP if email verification and result mail are required.",
    "Configure `AWS_REGION` and `S3_BUCKET` if S3-backed documents are required."
], left=0.8, top=1.75, width=8.2, height=5.2, font_size=18)


# Slide 11
slide = new_slide("EC2 and Production Deployment Steps", "How the project was promoted to cloud deployment")
add_bullets(slide, [
    "Create Ubuntu EC2 instance and allow SSH, HTTP, and optional HTTPS.",
    "Install Node.js, npm, PM2, Nginx, Git, and AWS CLI on the server.",
    "Attach IAM role to EC2 so the backend can access S3 securely.",
    "Clone the repository on EC2 and configure backend environment variables.",
    "Run backend with PM2.",
    "Build the frontend with Vite and copy `dist/` to the Nginx web root.",
    "Reload Nginx and verify API, login, and document access flow."
], left=0.8, top=1.75, width=8.2, height=5.3, font_size=18)


# Slide 12
slide = new_slide("Conclusion", "Final takeaway")
add_bullets(slide, [
    "The project combines academic workflow and integrity workflow in one platform.",
    "It uses a practical full-stack architecture with React, Express, PostgreSQL, EC2, and S3.",
    "The system supports secure exam delivery, manual review, controlled publication, and audit-ready reporting."
], left=0.95, top=2.0, width=7.8, height=2.0, font_size=22)
quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.5), Inches(7.8), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = NAVY
quote.line.fill.background()
p = quote.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Exam Integrity System = Online Exam Platform + Integrity Monitoring + Audit Reporting"
r.font.name = "Aptos"
r.font.bold = True
r.font.size = Pt(22)
r.font.color.rgb = WHITE


OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT))
print(f"Presentation created: {OUTPUT}")
